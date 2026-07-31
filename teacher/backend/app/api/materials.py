"""
教学资料与题库管理 API — PDF上传、AI出题、作业发布。
"""

from __future__ import annotations

import json
import os
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any

from fastapi import APIRouter, Depends, Form, HTTPException, UploadFile
from sqlalchemy.orm import Session
from app.core.llm import chat_json
from app.core.curriculum import get_course_list, get_chapters
from app.models.database import LessonPlan, get_db
from app.models.schemas import APIResponse

router = APIRouter(prefix="/api/materials", tags=["教学资料"])

# 存储路径
MATERIALS_DIR = Path(__file__).parent.parent.parent / "knowledge_base" / "materials"
QUESTIONS_DIR = Path(__file__).parent.parent.parent / "knowledge_base" / "questions"
MATERIALS_DIR.mkdir(parents=True, exist_ok=True)
QUESTIONS_DIR.mkdir(parents=True, exist_ok=True)

# 内存索引（生产环境应使用数据库）
_materials_index: dict[str, dict] = {}
_questions_index: dict[str, dict] = {}


def _load_indexes():
    """加载持久化索引。"""
    idx_file = MATERIALS_DIR / "_index.json"
    if idx_file.exists():
        try:
            data = json.loads(idx_file.read_text(encoding="utf-8"))
            _materials_index.update(data.get("materials", {}))
            _questions_index.update(data.get("questions", {}))
        except Exception:
            pass


def _save_indexes():
    """保存索引到磁盘。"""
    idx_file = MATERIALS_DIR / "_index.json"
    idx_file.write_text(
        json.dumps({"materials": _materials_index, "questions": _questions_index}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


# 启动时加载
_load_indexes()


# ── 资料管理 ──────────────────────────────────────────────

# AI 课程/章节识别提示词
COURSE_DETECT_PROMPT = """你是一个教学资料分析助手。根据提供的教材文本片段，判断它属于哪门课程、哪个章节。

可选的课程列表：{courses}

每门课程的章节列表：
{curriculum_summary}

请分析文本内容，返回 JSON：
{{
  "course": "最匹配的课程名（必须从可选课程列表中选一个，若无法判断则填 未分类）",
  "chapter": "最匹配的章节名（从该课程的章节中选，若无法判断则填空字符串）",
  "confidence": "高/中/低"
}}

注意：只返回 JSON，不要其他内容。"""


async def _detect_course_chapter(text: str) -> tuple[str, str]:
    """使用 AI 分析文本，自动识别课程和章节。"""
    courses = get_course_list()
    if not courses or len(text.strip()) < 100:
        return "未分类", ""

    # 构建课程体系摘要
    summary_parts = []
    for c in courses:
        chs = get_chapters(c)
        summary_parts.append(f"  {c}: {', '.join(chs[:5])}{'...' if len(chs) > 5 else ''}")
    curriculum_summary = "\n".join(summary_parts)

    try:
        result = chat_json(
            messages=[
                {"role": "system", "content": COURSE_DETECT_PROMPT.format(
                    courses=", ".join(courses),
                    curriculum_summary=curriculum_summary,
                )},
                {"role": "user", "content": f"请分析以下教材文本（前2000字），判断它属于哪门课程和章节：\n\n{text[:2000]}"},
            ],
            temperature=0.2,
        )
        course = result.get("course", "未分类")
        chapter = result.get("chapter", "")
        # 验证课程名是否在列表中
        if course not in courses and course != "未分类":
            course = "未分类"
        return course, chapter
    except Exception:
        return "未分类", ""


@router.post("/upload", response_model=APIResponse)
async def upload_material(file: UploadFile, course: str = Form(""), chapter: str = Form("")):
    """上传单个教学资料（PDF / Word）。课程和章节可选，不填则 AI 自动识别。"""
    return await _upload_single(file, course, chapter)


@router.post("/upload-batch", response_model=APIResponse)
async def upload_materials_batch(files: list[UploadFile], course: str = Form(""), chapter: str = Form("")):
    """批量上传多个教学资料。课程和章节可选，不填则 AI 自动识别。"""
    results = []
    success_count = 0
    for file in files:
        try:
            res = await _upload_single(file, course, chapter)
            if res.success:
                success_count += 1
                results.append({"filename": file.filename, "status": "success", "id": res.data["id"], "course": res.data["course"]})
            else:
                results.append({"filename": file.filename, "status": "failed", "message": res.message})
        except Exception as e:
            results.append({"filename": file.filename if file.filename else "unknown", "status": "error", "message": str(e)[:200]})

    return APIResponse(
        success=success_count > 0,
        message=f"成功上传 {success_count}/{len(files)} 个文件",
        data={"total": len(files), "success_count": success_count, "results": results},
    )


async def _upload_single(file: UploadFile, course: str, chapter: str) -> APIResponse:
    """处理单个文件上传。"""
    if not file.filename:
        return APIResponse(success=False, message="请选择文件")

    ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else ""
    if ext not in ("pdf", "docx", "doc", "pptx", "ppt", "zip", "png", "jpg", "jpeg"):
        return APIResponse(success=False, message=f"不支持的文件格式: .{ext}（支持 PDF、Word、PPT、图片、ZIP）")

    material_id = str(uuid.uuid4())[:8]
    save_path = MATERIALS_DIR / f"{material_id}_{file.filename}"
    content = await file.read()
    save_path.write_bytes(content)

    # 提取文本
    text_content = ""
    pages = 0
    try:
        if ext == "pdf":
            from pypdf import PdfReader
            import io
            reader = PdfReader(io.BytesIO(content))
            pages = len(reader.pages)
            text_content = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif ext in ("docx", "doc"):
            try:
                from docx import Document
                import io as _io
                doc = Document(_io.BytesIO(content))
                text_content = "\n".join(p.text for p in doc.paragraphs)
                pages = max(1, len(doc.paragraphs) // 40)
            except Exception:
                return APIResponse(success=False, message="无法解析 .doc 文件，请转换为 .docx 格式")
        elif ext in ("pptx", "ppt"):
            try:
                from pptx import Presentation
                import io as _io
                prs = Presentation(_io.BytesIO(content))
                slides_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slides_text.append(shape.text)
                text_content = "\n".join(slides_text)
                pages = len(prs.slides)
            except Exception:
                text_content = "（PPT 文本提取失败，文件已保存）"
                pages = 1
        elif ext in ("png", "jpg", "jpeg"):
            text_content = "（图片文件，已保存）"
            pages = 1
        elif ext == "zip":
            text_content = "（压缩包，已保存）"
            pages = 1
    except Exception as e:
        text_content = f"（文本提取失败: {str(e)}）"

    # AI 自动识别课程和章节（仅在未手动指定时）
    detected_course = course or ""
    detected_chapter = chapter or ""
    if (not course or not chapter) and text_content and len(text_content.strip()) >= 100:
        ai_course, ai_chapter = await _detect_course_chapter(text_content)
        if not course:
            detected_course = ai_course
        if not chapter:
            detected_chapter = ai_chapter

    info = {
        "id": material_id,
        "filename": file.filename,
        "course": detected_course or "未分类",
        "chapter": detected_chapter or "",
        "size": len(content),
        "size_display": f"{len(content) / 1024:.1f} KB" if len(content) < 1024 * 1024 else f"{len(content) / 1024 / 1024:.1f} MB",
        "pages": pages,
        "text_preview": text_content[:5000],
        "_source": "user",
	    "text_content": text_content,
        "created_at": datetime.now().isoformat()[:19],
    }
    _materials_index[material_id] = info
    _save_indexes()

    return APIResponse(success=True, message=f"上传成功：{file.filename}", data={
        "id": material_id,
        "filename": file.filename,
        "course": info["course"],
        "chapter": info["chapter"],
        "size_display": info["size_display"],
        "pages": info["pages"],
        "created_at": info["created_at"],
    })


@router.get("/list", response_model=APIResponse)
async def list_materials(course: str = ""):
    """获取教学资料列表。"""
    # 浅拷贝避免修改内存索引中的原始数据
    items = [dict(i) for i in _materials_index.values()]
    if course:
        items = [i for i in items if i["course"] == course]
    # 按时间倒序
    items.sort(key=lambda x: x.get("created_at", ""), reverse=True)
    # 不返回完整文本内容（仅预览），补齐 _source 标记
    for item in items:
        item.pop("text_content", None)
        if "_source" not in item:
            # 已知的系统种子文件列表
            seed_files = {"《神经网络与深度学习》.pdf", "大模型微调评估.pdf", "期末课程设计报告.pdf"}
            is_seed = item.get("id", "").startswith("seed_") or item.get("filename", "") in seed_files
            item["_source"] = "seed" if is_seed else "user"
    return APIResponse(success=True, data={"total": len(items), "items": items})


@router.get("/detail/{material_id}", response_model=APIResponse)
async def get_material(material_id: str):
    """获取单个资料详情。"""
    info = _materials_index.get(material_id)
    if not info:
        raise HTTPException(status_code=404, detail="资料不存在")
    return APIResponse(success=True, data=info)


_MIME_TYPES = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "doc": "application/msword",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "ppt": "application/vnd.ms-powerpoint",
    "zip": "application/zip",
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
}


@router.get("/download/{material_id}")
async def download_material(material_id: str):
    """下载已上传的教学资料文件。"""
    from fastapi.responses import FileResponse
    info = _materials_index.get(material_id)
    if not info:
        raise HTTPException(status_code=404, detail="资料不存在")
    for f in MATERIALS_DIR.iterdir():
        if f.name.startswith(material_id):
            ext = info["filename"].rsplit(".", 1)[-1].lower() if "." in info["filename"] else ""
            mime = _MIME_TYPES.get(ext, "application/octet-stream")
            return FileResponse(
                path=str(f),
                filename=info["filename"],
                media_type=mime,
            )
    raise HTTPException(status_code=404, detail="文件不存在，可能已被删除")


@router.post("/{material_id}/to-lesson", response_model=APIResponse)
async def material_to_lesson(material_id: str, db=Depends(get_db)):
    """将教学资料插入备课 → 创建教案记录，出现在教学台账中心。"""
    info = _materials_index.get(material_id)
    if not info:
        raise HTTPException(status_code=404, detail="资料不存在")

    # 提取文本内容作为教案的 sessions
    text = info.get("text_content", info.get("text_preview", ""))
    course_name = info.get("course", "未分类")
    chapter = info.get("chapter", "")
    filename = info.get("filename", "")

    # 将文本分块作为教学流程
    sessions = []
    paragraphs = [p.strip() for p in text.split("\n") if p.strip() and len(p.strip()) > 10]
    chunk_size = max(1, len(paragraphs) // 4) if len(paragraphs) >= 4 else len(paragraphs)
    for i in range(0, min(len(paragraphs), chunk_size * 4), chunk_size):
        chunk = paragraphs[i:i + chunk_size]
        if chunk:
            sessions.append({
                "session_order": len(sessions) + 1,
                "session_topic": f"第{len(sessions) + 1}部分：{filename}",
                "key_points": [chunk[0][:30]] if chunk else [],
                "activities": [{
                    "activity_type": "讲授",
                    "duration": 15,
                    "content": "\n".join(chunk[:5]),
                }],
            })

    plan_data = {
        "course_name": course_name,
        "chapter": chapter or "教材内容",
        "total_hours": max(1, len(sessions)),
        "objectives": [{"dimension": "知识目标", "content": f"掌握「{filename}」中的核心知识点"}],
        "methods": ["讲授法", "讨论法"],
        "resources": ["教材", "多媒体课件"],
        "sessions": sessions,
        "board_design": {},
        "class_tasks": [],
        "homework": [],
        "assessment": {},
        "innovation": {},
    }

    lesson = LessonPlan(
        id=f"mat_{material_id}",
        course_name=course_name,
        chapter=chapter or "教材内容",
        total_hours=max(1, len(sessions)),
        additional_requirements=f"从资源中心导入：{filename}",
        plan_data=json.dumps(plan_data, ensure_ascii=False),
        created_at=datetime.now(),
        updated_at=datetime.now(),
    )
    db.add(lesson)
    db.commit()

    # 审计日志
    from app.api.audit import log_operation, save_snapshot
    log_operation(db, lesson.id, "create", operator="教师",
                  course_name=course_name, chapter=chapter or "教材内容",
                  detail=f"从资源中心插入备课：{filename}",
                  plan_name=f"{course_name} — {chapter or '教材内容'}")
    save_snapshot(db, lesson.id, plan_data, created_by="教师")

    return APIResponse(success=True, message=f"「{filename}」已插入备课", data={"plan_id": lesson.id})


@router.delete("/delete/{material_id}", response_model=APIResponse)
async def delete_material(material_id: str):
    """删除教学资料，同时清理关联的题目和发布通知。"""
    info = _materials_index.pop(material_id, None)
    if not info:
        raise HTTPException(status_code=404, detail="资料不存在")
    # 删除文件
    for f in MATERIALS_DIR.iterdir():
        if f.name.startswith(material_id):
            f.unlink(missing_ok=True)

    # 清理关联的题目
    removed_qids = [qid for qid, q in list(_questions_index.items()) if q.get("material_id") == material_id]
    for qid in removed_qids:
        _questions_index.pop(qid, None)

    # 清理关联的发布记录及其通知
    for f in QUESTIONS_DIR.glob("publish_*.json"):
        try:
            record = json.loads(f.read_text(encoding="utf-8"))
            pub_qids = record.get("question_ids", [])
            # 如果发布记录中的所有题目都属于被删的资料，则清理该发布
            if pub_qids and all(qid in removed_qids for qid in pub_qids):
                notif_id = record.get("notification_id")
                if notif_id:
                    try:
                        from app.api.notifications import _read, _write as _write_notifs
                        notifs = _read()
                        notifs = [n for n in notifs if n.get("id") != notif_id]
                        _write_notifs(notifs)
                    except Exception:
                        pass
                f.unlink(missing_ok=True)
        except Exception:
            continue

    _save_indexes()
    return APIResponse(success=True, message=f"已删除资料及关联的 {len(removed_qids)} 道题目")


# ── AI 出题 ────────────────────────────────────────────────

QUESTION_SYSTEM_PROMPT = """【当前任务：学科专业试题智能出题与分层命题】
你是一流学科建设专家型教师，根据教材内容生成标准化本科试题。

出题要求：
1. 分层出题：基础题、提高题、综合应用题、前沿创新题
2. 题型包含：选择、判断、简答、计算、案例分析、论述（按需适配）
3. 每题包含：题目、标准答案、分步评分细则、详细解析、易错点分析
4. 每题绑定知识点标签、教学目标、命题依据（教材来源/页码）
5. 试题规避网络原题，具备本科专业高阶考察性
6. 难度分布合理，覆盖不同认知层次（记忆/理解/应用/分析/评价/创造）

【数学符号强制规范】所有数学表达式必须使用标准符号：
❌ "x的n次方求和" → ✅ Σₙ₌₀ᐁ xⁿ
❌ "a_n乘以x的n次方" → ✅ aₙxⁿ 或 ∑ aₙxⁿ
❌ "f对x求导" → ✅ df/dx 或 f'(x)
❌ "从a到b的积分" → ✅ ∫ₐᵇ
用Unicode数学符号：∑ ∫ ∂ √ ² ³ ⁿ ∞ → ∈ ⊂ α β γ Δ π ≤ ≥ ≠ ≈ ± × · ₁ ₂ ₙ ½ ⅓ ⅔ ¼ ¾
分数用对角线形式：a/b、(x+1)/(x−1)，严禁用中文如"a分之b""二分之一"。
复杂分式必须加括号：(x²+1)/(x−1) ≠ x²+1/x−1

末尾添加 AI 生成标识。

输出 JSON 格式：
{
  "questions": [
    {
      "question": "题目内容",
      "type": "选择题|填空题|简答题|计算题|论述题|案例分析",
      "options": ["A. xxx", "B. xxx", "C. xxx", "D. xxx"],
      "answer": "标准答案",
      "difficulty": "基础|提高|综合|前沿创新",
      "knowledge_point": "知识点名称",
      "teaching_objective": "对应教学目标",
      "source": "命题依据（教材章节/页码/文献）",
      "scoring_rubric": "分步评分细则",
      "common_mistakes": "学生常见易错点",
      "explanation": "详细解析及解题思路",
      "cognitive_level": "记忆|理解|应用|分析|评价|创造",
      "estimated_time": 5
    }
  ]
}"""


@router.post("/generate-questions", response_model=APIResponse)
async def generate_questions(data: dict):
    """
    基于教学资料生成练习题。

    请求体：
    {
        "material_id": "xxx",
        "count": 5,
        "difficulty": "中等",
        "types": ["选择题", "填空题", "简答题"]
    }
    """
    material_id = data.get("material_id", "")
    count = min(data.get("count", 5), 20)
    difficulty = data.get("difficulty", "中等")
    types = data.get("types", ["选择题", "填空题", "简答题"])
    question_types = data.get("question_types", None) or types  # 兼容两种参数名

    info = _materials_index.get(material_id)
    if not info:
        raise HTTPException(status_code=404, detail="资料不存在，请先上传PDF")

    text = info.get("text_content", "")
    if not text or len(text.strip()) < 50:
        raise HTTPException(status_code=400, detail="资料文本内容不足，无法出题")

    # 截取前 8000 字符作为上下文
    context = text[:8000]

    user_prompt = f"""请根据以下教材内容，生成 {count} 道{difficulty}难度的练习题。

教材内容（{info['filename']} - {info['chapter'] or info['course']}）：
{context}

题目类型：{', '.join(question_types)}
难度级别：{difficulty}
题目数量：{count} 题

请确保覆盖不同知识点，难度分布合理。"""

    try:
        result = chat_json(
            messages=[
                {"role": "system", "content": QUESTION_SYSTEM_PROMPT},
                {"role": "user", "content": user_prompt},
            ],
            temperature=0.5,
        )
        questions = result.get("questions", [])
        if not questions and "exercises" in result:
            questions = result["exercises"]
    except Exception as e:
        err_msg = str(e)
        if "401" in err_msg or "auth" in err_msg.lower() or "credential" in err_msg.lower():
            return APIResponse(success=False, message="LLM API Key 无效或已过期，请前往「LLM API 配置」页面更新密钥")
        if "connection" in err_msg.lower() or "timeout" in err_msg.lower():
            return APIResponse(success=False, message="LLM 服务连接失败，请检查网络或 API 地址配置")
        return APIResponse(success=False, message=f"出题失败：{err_msg[:100]}")

    # 保存题目
    batch_id = str(uuid.uuid4())[:8]
    saved_questions = []
    for i, q in enumerate(questions):
        qid = f"{batch_id}_{i}"
        question_item = {
            "id": qid,
            "batch_id": batch_id,
            "material_id": material_id,
            "material_name": info["filename"],
            "course": info["course"],
            "question": q.get("question", ""),
            "type": q.get("type", "简答题"),
            "options": q.get("options", []),
            "answer": q.get("answer", ""),
            "difficulty": q.get("difficulty", "中等"),
            "knowledge_point": q.get("knowledge_point", ""),
            "explanation": q.get("explanation", ""),
            "estimated_time": q.get("estimated_time", 5),
            "status": "draft",  # draft / published
            "created_at": datetime.now().isoformat()[:19],
        }
        _questions_index[qid] = question_item
        saved_questions.append(question_item)

    # 按题型、难度升序排列
    type_order = {"选择题": 1, "填空题": 2, "简答题": 3, "计算题": 4, "论述题": 5, "案例分析": 6}
    diff_order = {"基础": 1, "提高": 2, "中等": 3, "综合": 4, "前沿创新": 5, "前沿": 5}
    saved_questions.sort(key=lambda q: (
        type_order.get(q.get("type", ""), 99),
        diff_order.get(q.get("difficulty", ""), 99),
    ))

    _save_indexes()

    return APIResponse(
        success=True,
        message=f"成功生成 {len(saved_questions)} 道题目",
        data={
            "batch_id": batch_id,
            "total": len(saved_questions),
            "questions": saved_questions,
        },
    )


@router.get("/questions", response_model=APIResponse)
async def list_questions(material_id: str = "", status: str = "", course: str = ""):
    """获取题目列表。"""
    items = list(_questions_index.values())
    if material_id:
        items = [i for i in items if i["material_id"] == material_id]
    if status:
        items = [i for i in items if i["status"] == status]
    if course:
        items = [i for i in items if i.get("course") == course]
    items.sort(key=lambda x: x.get("created_at", ""), reverse=True)
    return APIResponse(success=True, data={"total": len(items), "items": items})


@router.post("/questions/update", response_model=APIResponse)
async def update_question(data: dict):
    """更新单道题目。"""
    qid = data.get("id", "")
    if qid not in _questions_index:
        raise HTTPException(status_code=404, detail="题目不存在")
    for key in ("question", "answer", "explanation", "difficulty", "type", "options", "knowledge_point"):
        if key in data:
            _questions_index[qid][key] = data[key]
    _save_indexes()
    return APIResponse(success=True, message="已更新")


@router.post("/publish", response_model=APIResponse)
async def publish_questions(data: dict):
    """
    发布题目给学生端。

    请求体：
    {
        "question_ids": ["id1", "id2"],
        "course": "离散数学",
        "title": "第三章课后练习",
        "deadline": "2026-06-20"
    }
    """
    qids = data.get("question_ids", [])
    course = data.get("course", "")
    title = data.get("title", "练习题")
    deadline = data.get("deadline", "")

    published = []
    for qid in qids:
        if qid in _questions_index:
            _questions_index[qid]["status"] = "published"
            published.append(qid)

    if not published:
        raise HTTPException(status_code=400, detail="没有可发布的题目")

    _save_indexes()

    # 先创建通知，以便将通知ID存入发布记录（撤销时可联动删除）
    notification_id = None
    try:
        from app.api.notifications import create_notification
        deadline_info = f"，截止 {deadline}" if deadline else ""
        notif = create_notification(
            ntype="作业",
            title=f"{course} · {title}",
            desc=f"已发布 {len(published)} 道题目{deadline_info}",
            route="/materials",
        )
        notification_id = notif.get("id")
    except Exception:
        pass  # 通知创建失败不影响发布

    # 构建发布记录（含通知ID）
    publish_record = {
        "id": str(uuid.uuid4())[:8],
        "title": title,
        "course": course,
        "deadline": deadline,
        "question_count": len(published),
        "question_ids": published,
        "notification_id": notification_id,
        "created_at": datetime.now().isoformat()[:19],
    }

    pub_file = QUESTIONS_DIR / f"publish_{publish_record['id']}.json"
    pub_file.write_text(json.dumps(publish_record, ensure_ascii=False, indent=2), encoding="utf-8")

    return APIResponse(
        success=True,
        message=f"已成功发布 {len(published)} 道题目",
        data=publish_record,
    )


@router.get("/publish/list", response_model=APIResponse)
async def list_published():
    """获取已发布的作业列表。"""
    records = []
    for f in QUESTIONS_DIR.glob("publish_*.json"):
        try:
            records.append(json.loads(f.read_text(encoding="utf-8")))
        except Exception:
            continue
    records.sort(key=lambda x: x.get("created_at", ""), reverse=True)
    return APIResponse(success=True, data={"total": len(records), "items": records})


@router.get("/publish/{publish_id}/questions", response_model=APIResponse)
async def get_published_questions(publish_id: str):
    """获取某次发布的完整题目列表。"""
    pub_file = QUESTIONS_DIR / f"publish_{publish_id}.json"
    if not pub_file.exists():
        raise HTTPException(status_code=404, detail="发布记录不存在")
    record = json.loads(pub_file.read_text(encoding="utf-8"))
    questions = []
    for qid in record.get("question_ids", []):
        if qid in _questions_index:
            questions.append(_questions_index[qid])
    return APIResponse(success=True, data={"title": record.get("title", ""), "course": record.get("course", ""), "deadline": record.get("deadline", ""), "questions": questions, "total": len(questions)})


@router.post("/questions/export-word", response_model=APIResponse)
async def export_questions_word(data: dict):
    """将题目列表导出为 Word 文档（返回 base64 内容）。"""
    qids = data.get("question_ids", [])
    title = data.get("title", "习题集")
    questions = []
    for qid in qids:
        if qid in _questions_index:
            questions.append(_questions_index[qid])
    if not questions:
        raise HTTPException(status_code=400, detail="没有可导出的题目")

    # 构建 Word 文档内容（纯文本转 Word）
    lines = [f"《{title}》", "", f"共 {len(questions)} 道题", "=" * 40, ""]
    for idx, q in enumerate(questions, 1):
        lines.append(f"{idx}. [{q.get('type', '')}] [{q.get('difficulty', '')}] {q.get('question', '')}")
        opts = q.get("options", [])
        if opts:
            for opt in opts:
                lines.append(f"   {opt}")
        lines.append(f"   答案：{q.get('answer', '')}")
        if q.get("explanation"):
            lines.append(f"   解析：{q.get('explanation', '')}")
        if q.get("knowledge_point"):
            lines.append(f"   知识点：{q.get('knowledge_point', '')}")
        lines.append("")

    try:
        from app.api.resources import _build_docx
        doc_bytes = _build_docx(title=title, body="\n".join(lines))
        import base64
        b64 = base64.b64encode(doc_bytes).decode()
        return APIResponse(success=True, message="导出成功", data={"base64": b64, "filename": f"{title}.docx"})
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Word 导出失败：{str(e)[:100]}")


@router.post("/save-exercises", response_model=APIResponse)
async def save_exercises_to_bank(data: dict):
    """
    将出题助手生成的练习题保存到题库（_questions_index）。

    请求体：
    {
        "course_name": "机器学习",
        "chapter": "第三章",
        "exercises": [
            {"question": "...", "type": "选择题", "options": ["A.xx","B.xx"], "answer": "B", ...},
            ...
        ]
    }
    """
    course = data.get("course_name", "")
    chapter = data.get("chapter", "")
    exercises = data.get("exercises", [])
    if not exercises:
        raise HTTPException(status_code=400, detail="没有可保存的题目")

    batch_id = str(uuid.uuid4())[:8]
    saved = []
    for i, ex in enumerate(exercises):
        qid = f"ex_{batch_id}_{i}"
        question_item = {
            "id": qid,
            "batch_id": batch_id,
            "material_id": "",
            "material_name": f"{course} · {chapter or '综合练习'}",
            "course": course,
            "chapter": chapter,
            "question": ex.get("question", ""),
            "type": ex.get("type", "简答题"),
            "options": ex.get("options", []),
            "answer": ex.get("answer", ""),
            "difficulty": ex.get("difficulty", "中等"),
            "knowledge_point": ex.get("knowledge_point", ""),
            "explanation": ex.get("explanation", ""),
            "estimated_time": ex.get("estimated_time", 5),
            "status": "draft",
            "created_at": datetime.now().isoformat()[:19],
        }
        _questions_index[qid] = question_item
        saved.append(question_item)

    _save_indexes()

    return APIResponse(
        success=True,
        message=f"已保存 {len(saved)} 道题目到题库",
        data={"batch_id": batch_id, "total": len(saved), "questions": saved},
    )


@router.post("/questions/unpublish", response_model=APIResponse)
async def unpublish_questions(data: dict):
    """
    撤销发布：将指定发布记录中的题目恢复为草稿状态，删除发布记录。
    """
    publish_id = data.get("publish_id", "")
    if not publish_id:
        raise HTTPException(status_code=400, detail="缺少发布记录ID")

    pub_file = QUESTIONS_DIR / f"publish_{publish_id}.json"
    if not pub_file.exists():
        raise HTTPException(status_code=404, detail="发布记录不存在")

    record = json.loads(pub_file.read_text(encoding="utf-8"))
    qids = record.get("question_ids", [])
    count = 0
    for qid in qids:
        if qid in _questions_index:
            _questions_index[qid]["status"] = "draft"
            count += 1

    pub_file.unlink(missing_ok=True)
    _save_indexes()

    # 联动删除发布时创建的通知
    notification_id = record.get("notification_id")
    if notification_id:
        try:
            from app.api.notifications import _read, _write as _write_notifs
            notifs = _read()
            before = len(notifs)
            notifs = [n for n in notifs if n.get("id") != notification_id]
            if len(notifs) < before:
                _write_notifs(notifs)
        except Exception:
            pass  # 通知删除失败不影响撤销

    return APIResponse(success=True, message=f"已撤销 {count} 道题目，恢复为草稿")


@router.post("/questions/clear-orphaned", response_model=APIResponse)
async def clear_orphaned_questions(data: dict = {}):
    """清空孤立题目或全部草稿题目。

    请求体（可选）：
        all_drafts: bool — 为 True 时清空全部草稿题目，否则仅清空孤立题目
    """
    clear_all = data.get("all_drafts", False) if data else False

    if clear_all:
        # 清空全部草稿题目（保留已发布）
        removed = [qid for qid, q in list(_questions_index.items())
                   if q.get("status") != "published"]
        for qid in removed:
            _questions_index.pop(qid, None)
        _save_indexes()
        return APIResponse(
            success=True,
            message=f"已清空 {len(removed)} 道草稿题目",
            data={"cleared": len(removed)},
        )
    else:
        # 仅清空孤立题目
        valid_material_ids = set(_materials_index.keys())
        orphaned = [qid for qid, q in list(_questions_index.items())
                    if q.get("material_id") and q["material_id"] not in valid_material_ids]
        for qid in orphaned:
            _questions_index.pop(qid, None)
        if orphaned:
            _save_indexes()
        return APIResponse(
            success=True,
            message=f"已清理 {len(orphaned)} 道孤立题目",
            data={"cleared": len(orphaned)},
        )
